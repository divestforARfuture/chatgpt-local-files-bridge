from __future__ import annotations

from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .errors import UnsupportedFileError

_TEXT_SUFFIXES = {
    ".bat", ".c", ".cfg", ".cmd", ".conf", ".cpp", ".css", ".csv",
    ".go", ".h", ".hpp", ".html", ".ini", ".java", ".js", ".json",
    ".log", ".md", ".ps1", ".py", ".rb", ".rs", ".sh", ".sql",
    ".toml", ".ts", ".txt", ".xml", ".yaml", ".yml",
}


def _read_plain_text(path: Path) -> str:
    data = path.read_bytes()
    for encoding in ("utf-8-sig", "utf-8", "cp1252"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise UnsupportedFileError("The file is not valid text in a supported encoding.")


def _read_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append(f"\n--- Page {index} ---\n{page.extract_text() or ''}")
    return "".join(pages).lstrip()


def _read_docx(path: Path) -> str:
    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            paragraphs.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(paragraphs)


def _read_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    chunks: list[str] = []
    try:
        for sheet in workbook.worksheets:
            chunks.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                chunks.append("\t".join("" if value is None else str(value) for value in row))
    finally:
        workbook.close()
    return "\n".join(chunks)


def _read_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    chunks: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        chunks.append(f"--- Slide {index} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def extract_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.casefold()
    if suffix in _TEXT_SUFFIXES or suffix == "":
        return "text", _read_plain_text(path)
    if suffix == ".pdf":
        return "pdf", _read_pdf(path)
    if suffix == ".docx":
        return "docx", _read_docx(path)
    if suffix == ".xlsx":
        return "xlsx", _read_xlsx(path)
    if suffix == ".pptx":
        return "pptx", _read_pptx(path)
    raise UnsupportedFileError(
        "Unsupported file type.",
        details={"suffix": suffix or "(none)"},
    )
